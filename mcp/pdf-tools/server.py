#!/usr/bin/env python3
"""MAG Document Tools MCP server (stdio, Python).

Exposes document utility tools to the Hermes agent that are unavailable via
execute_code on client channels (WhatsApp/Telegram). Only reads/writes files
within /opt/data — no arbitrary filesystem access. Started as PDF-only; now
covers the other office formats a client can attach in chat (Word/Excel/
PowerPoint/plain text), since they hit the exact same execute_code-unavailable
wall PDFs did.

Tools:
  extract_pdf_text      — extract plain text from a cached PDF (page-range pagination,
                           OCR fallback for scanned pages, capped OCR page count)
  extract_pdf_images    — extract embedded JPEG/PNG images from a cached PDF
  extract_docx_text     — extract paragraphs+tables text from a cached Word document
  extract_xlsx_text     — extract sheet contents (CSV-ish) from a cached Excel workbook
  extract_pptx_text     — extract slide text+notes from a cached PowerPoint deck
  read_text_file        — read a cached plain-text file (.txt/.csv/.json/.md/...)
  generate_pdf_report   — create a PDF report with title, date, body text, images and captions
  pdf_split              — extract a page range from a PDF into a new PDF
  pdf_merge               — concatenate multiple PDFs into one
  pdf_watermark           — stamp a diagonal text watermark on every page of a PDF
"""

import json
import os
import sys

SERVER_NAME = "mag-pdf-tools"
SERVER_VERSION = "0.5.0"
PROTOCOL_VERSION = "2025-06-18"

# Wall-clock-cheap cap on returned text so a huge/many-page document can't blow
# the MCP response or downstream context. For PDFs, prefer start_page/end_page
# to page through a long document instead of raising this; for docx/text files
# (no natural page unit), raise it or use the offset/next_offset pair instead.
DEFAULT_MAX_TEXT_CHARS = 20000

# Hard cap on how many scanned (no text layer) pages get OCR'd in a single
# extract_pdf_text call. Tesseract at 300dpi can take several seconds a page —
# a bulky scanned document OCR'd unbounded risks stalling the whole tool call
# past the turn's budget. Mirrors extract_pdf_images' existing max_images cap.
# Pages beyond the cap are skipped (not silently dropped — reported back so the
# agent can re-call with a narrower start_page/end_page to reach them).
DEFAULT_MAX_OCR_PAGES = 15

# Plain-text formats read_text_file will open directly (no binary/office formats).
TEXT_FILE_EXTENSIONS = {
    ".txt", ".csv", ".tsv", ".json", ".md", ".markdown", ".yaml", ".yml",
    ".log", ".xml", ".html", ".htm", ".ini", ".conf",
}
# Read cap in bytes — small text attachments only; anything bigger almost
# certainly isn't meant to be read verbatim in a chat reply.
MAX_TEXT_FILE_BYTES = 2 * 1024 * 1024

# Row cap per Excel sheet — protects against a workbook with tens of thousands
# of rows blowing the response; the agent can target a specific sheet_name to
# read more of one sheet instead of dumping everything.
DEFAULT_MAX_ROWS_PER_SHEET = 500

# Longer side any embedded image is downscaled to before base64-embedding (keeps
# HTML size + chromium render time bounded for multi-photo reports, e.g. full-res
# inspection photos). Skipped gracefully if Pillow isn't importable.
MAX_IMAGE_DIMENSION = 1600

# Wall-clock budget for the chromium print-to-pdf subprocess. Bumped from the
# original 60s: several downscaled-but-still-numerous embedded photos can still
# take a while to lay out and rasterize under load.
CHROMIUM_TIMEOUT_SECONDS = 120

OUTPUT_DIR = "/opt/data/workspace"


def log(*args):
    print(f"[mag-pdf-tools] {' '.join(str(a) for a in args)}", file=sys.stderr, flush=True)


def send(obj):
    sys.stdout.write(json.dumps(obj) + "\n")
    sys.stdout.flush()


def reply(id_, result):
    send({"jsonrpc": "2.0", "id": id_, "result": result})


def reply_error(id_, code, msg):
    send({"jsonrpc": "2.0", "id": id_, "error": {"code": code, "message": msg}})


TOOLS = [
    {
        "name": "extract_pdf_text",
        "description": (
            "Extrai o texto de um PDF (relatórios, estudos, contratos, artigos) para você ler e "
            "responder sobre o conteúdo. Tenta o texto nativo do PDF primeiro; se uma página vier "
            "vazia (PDF escaneado/imagem), cai automaticamente para OCR (português+inglês) nela — "
            "até um limite de páginas por chamada (veja max_ocr_pages). "
            "Este é o único caminho para ler o TEXTO de um PDF em canais de cliente "
            "(WhatsApp/Telegram) — não depende de execute_code. Para PDFs longos, use start_page/"
            "end_page para paginar em vez de só aumentar max_chars — assim você chega às páginas "
            "finais em vez de sempre reler o começo. NÃO use para extrair fotos/figuras embutidas "
            "(use extract_pdf_images para isso)."
        ),
        "inputSchema": {
            "type": "object",
            "properties": {
                "pdf_path": {
                    "type": "string",
                    "description": (
                        "Caminho absoluto para o arquivo PDF. "
                        "Normalmente: /opt/data/cache/documents/doc_<hash>_<nome>.pdf"
                    ),
                },
                "start_page": {
                    "type": "integer",
                    "description": "Primeira página a ler (1-indexado, padrão: 1). Use para paginar um PDF longo.",
                },
                "end_page": {
                    "type": "integer",
                    "description": "Última página a ler, inclusive (padrão: última página do PDF).",
                },
                "max_chars": {
                    "type": "integer",
                    "description": "Máximo de caracteres a devolver (padrão: 20000). Reduza se só precisar de um resumo.",
                    "default": DEFAULT_MAX_TEXT_CHARS,
                },
                "max_ocr_pages": {
                    "type": "integer",
                    "description": (
                        "Máximo de páginas escaneadas a rodar OCR nesta chamada (padrão: 15) — protege contra "
                        "PDFs escaneados grandes travarem a chamada. Páginas além do limite são reportadas, "
                        "não silenciosamente ignoradas — refine start_page/end_page para alcançá-las."
                    ),
                    "default": DEFAULT_MAX_OCR_PAGES,
                },
            },
            "required": ["pdf_path"],
        },
    },
    {
        "name": "extract_pdf_images",
        "description": (
            "Extrai imagens embutidas de um PDF (fotos, figuras, gráficos) e salva como arquivos de imagem separados. "
            "Retorna os caminhos de cada imagem extraída — use MEDIA:<path> em sua resposta para enviá-las ao chat. "
            "Use quando o usuário pedir para mostrar fotos/imagens de dentro de um PDF."
        ),
        "inputSchema": {
            "type": "object",
            "properties": {
                "pdf_path": {
                    "type": "string",
                    "description": (
                        "Caminho absoluto para o arquivo PDF. "
                        "Normalmente: /opt/data/cache/documents/doc_<hash>_<nome>.pdf"
                    ),
                },
                "max_images": {
                    "type": "integer",
                    "description": "Número máximo de imagens a extrair (padrão: 10). Use valores menores para economizar.",
                    "default": 10,
                },
                "min_size": {
                    "type": "integer",
                    "description": "Largura mínima em pixels para incluir uma imagem (padrão: 100). Filtra ícones/logos pequenos.",
                    "default": 100,
                },
            },
            "required": ["pdf_path"],
        },
    },
    {
        "name": "extract_docx_text",
        "description": (
            "Extrai o texto de um documento Word (.docx) — parágrafos e tabelas — para você ler e "
            "responder sobre o conteúdo. Único caminho para ler um .docx em canais de cliente "
            "(WhatsApp/Telegram); não depende de execute_code."
        ),
        "inputSchema": {
            "type": "object",
            "properties": {
                "docx_path": {
                    "type": "string",
                    "description": (
                        "Caminho absoluto para o arquivo .docx. "
                        "Normalmente: /opt/data/cache/documents/doc_<hash>_<nome>.docx"
                    ),
                },
                "offset": {
                    "type": "integer",
                    "description": "Posição (em caracteres) de onde começar a devolver texto (padrão: 0). Use next_offset da resposta anterior para continuar um documento longo.",
                    "default": 0,
                },
                "max_chars": {
                    "type": "integer",
                    "description": "Máximo de caracteres a devolver a partir de offset (padrão: 20000).",
                    "default": DEFAULT_MAX_TEXT_CHARS,
                },
            },
            "required": ["docx_path"],
        },
    },
    {
        "name": "extract_xlsx_text",
        "description": (
            "Extrai o conteúdo de uma planilha Excel (.xlsx) — cada aba vira um bloco de texto tipo "
            "CSV — para você ler e responder sobre os dados. Único caminho para ler um .xlsx em "
            "canais de cliente (WhatsApp/Telegram); não depende de execute_code."
        ),
        "inputSchema": {
            "type": "object",
            "properties": {
                "xlsx_path": {
                    "type": "string",
                    "description": (
                        "Caminho absoluto para o arquivo .xlsx. "
                        "Normalmente: /opt/data/cache/documents/doc_<hash>_<nome>.xlsx"
                    ),
                },
                "sheet_name": {
                    "type": "string",
                    "description": "Nome de uma aba específica a ler (padrão: todas as abas). Use quando a planilha tiver muitas abas ou uma aba grande foi cortada.",
                },
                "max_chars": {
                    "type": "integer",
                    "description": "Máximo de caracteres a devolver (padrão: 20000).",
                    "default": DEFAULT_MAX_TEXT_CHARS,
                },
                "max_rows_per_sheet": {
                    "type": "integer",
                    "description": "Máximo de linhas por aba nesta chamada (padrão: 500). Use sheet_name para ler mais linhas de uma aba específica.",
                    "default": DEFAULT_MAX_ROWS_PER_SHEET,
                },
            },
            "required": ["xlsx_path"],
        },
    },
    {
        "name": "extract_pptx_text",
        "description": (
            "Extrai o texto de uma apresentação PowerPoint (.pptx) — título, corpo, tabelas e notas "
            "do apresentador de cada slide — para você ler e responder sobre o conteúdo. Único "
            "caminho para ler um .pptx em canais de cliente (WhatsApp/Telegram); não depende de "
            "execute_code."
        ),
        "inputSchema": {
            "type": "object",
            "properties": {
                "pptx_path": {
                    "type": "string",
                    "description": (
                        "Caminho absoluto para o arquivo .pptx. "
                        "Normalmente: /opt/data/cache/documents/doc_<hash>_<nome>.pptx"
                    ),
                },
                "start_slide": {
                    "type": "integer",
                    "description": "Primeiro slide a ler (1-indexado, padrão: 1).",
                },
                "end_slide": {
                    "type": "integer",
                    "description": "Último slide a ler, inclusive (padrão: último slide).",
                },
                "max_chars": {
                    "type": "integer",
                    "description": "Máximo de caracteres a devolver (padrão: 20000).",
                    "default": DEFAULT_MAX_TEXT_CHARS,
                },
            },
            "required": ["pptx_path"],
        },
    },
    {
        "name": "read_text_file",
        "description": (
            "Lê um arquivo de texto simples (.txt, .csv, .tsv, .json, .md, .yaml, .log, .xml, .html) "
            "enviado em chat. Use para qualquer anexo de texto puro que não seja PDF/Word/Excel/"
            "PowerPoint. Não use para arquivos binários (imagens, PDFs, .docx/.xlsx/.pptx — cada um "
            "tem sua própria ferramenta)."
        ),
        "inputSchema": {
            "type": "object",
            "properties": {
                "file_path": {
                    "type": "string",
                    "description": "Caminho absoluto para o arquivo de texto.",
                },
                "max_chars": {
                    "type": "integer",
                    "description": "Máximo de caracteres a devolver (padrão: 20000).",
                    "default": DEFAULT_MAX_TEXT_CHARS,
                },
            },
            "required": ["file_path"],
        },
    },
    {
        "name": "generate_pdf_report",
        "description": (
            "Gera um relatório em PDF com título, data, texto narrativo opcional e imagens com legendas. "
            "Use quando o usuário pedir para criar/gerar um documento PDF com fotos (com ou sem texto/laudo). "
            "Este é o único caminho de geração de PDF disponível em canais de cliente (WhatsApp/Telegram) — "
            "não depende de execute_code. Se você já extraiu texto do PDF original (ex: via pymupdf4llm), "
            "passe esse texto em 'body' em vez de tentar montar HTML/chromium manualmente. "
            "Retorna o caminho do PDF gerado — inclua MEDIA:<path> na resposta para enviá-lo."
        ),
        "inputSchema": {
            "type": "object",
            "properties": {
                "title": {
                    "type": "string",
                    "description": "Título do relatório.",
                },
                "body": {
                    "type": "string",
                    "description": (
                        "Texto narrativo do relatório (ex: laudo/achados extraídos do PDF original). "
                        "Texto simples com parágrafos separados por linha em branco — sem HTML. "
                        "Renderizado entre o subtítulo e a grade de imagens."
                    ),
                    "default": "",
                },
                "images": {
                    "type": "array",
                    "description": "Lista de imagens a incluir.",
                    "items": {
                        "type": "object",
                        "properties": {
                            "path": {"type": "string", "description": "Caminho absoluto da imagem."},
                            "caption": {"type": "string", "description": "Legenda da imagem."},
                        },
                        "required": ["path"],
                    },
                },
                "subtitle": {
                    "type": "string",
                    "description": "Subtítulo ou descrição opcional.",
                    "default": "",
                },
                "output_name": {
                    "type": "string",
                    "description": "Nome do arquivo de saída sem extensão (padrão: 'relatorio').",
                    "default": "relatorio",
                },
                "images_per_row": {
                    "type": "integer",
                    "description": "Imagens por linha no layout (1 ou 2, padrão: 1).",
                    "default": 1,
                },
            },
            "required": ["title"],
        },
    },
    {
        "name": "pdf_split",
        "description": (
            "Extrai um intervalo de páginas de um PDF para um novo arquivo PDF. Use quando o "
            "usuário pedir para separar/recortar páginas específicas de um PDF. Único caminho "
            "disponível em canais de cliente (WhatsApp/Telegram) — não depende de execute_code."
        ),
        "inputSchema": {
            "type": "object",
            "properties": {
                "pdf_path": {"type": "string", "description": "Caminho absoluto do PDF de origem."},
                "start_page": {"type": "integer", "description": "Primeira página a incluir (1-indexado)."},
                "end_page": {"type": "integer", "description": "Última página a incluir, inclusive."},
                "output_name": {
                    "type": "string",
                    "description": "Nome do arquivo de saída sem extensão (padrão: 'split').",
                    "default": "split",
                },
            },
            "required": ["pdf_path", "start_page", "end_page"],
        },
    },
    {
        "name": "pdf_merge",
        "description": (
            "Concatena vários PDFs (na ordem dada) em um único arquivo PDF. Use quando o usuário "
            "pedir para juntar/mesclar PDFs. Único caminho disponível em canais de cliente "
            "(WhatsApp/Telegram) — não depende de execute_code."
        ),
        "inputSchema": {
            "type": "object",
            "properties": {
                "pdf_paths": {
                    "type": "array",
                    "description": "Caminhos absolutos dos PDFs a mesclar, na ordem desejada (mínimo 2).",
                    "items": {"type": "string"},
                },
                "output_name": {
                    "type": "string",
                    "description": "Nome do arquivo de saída sem extensão (padrão: 'merged').",
                    "default": "merged",
                },
            },
            "required": ["pdf_paths"],
        },
    },
    {
        "name": "pdf_watermark",
        "description": (
            "Aplica uma marca d'água de texto diagonal em todas as páginas de um PDF. Use quando o "
            "usuário pedir para marcar/carimbar um PDF (ex: 'CONFIDENCIAL', 'RASCUNHO', nome da "
            "empresa). Único caminho disponível em canais de cliente (WhatsApp/Telegram) — não "
            "depende de execute_code."
        ),
        "inputSchema": {
            "type": "object",
            "properties": {
                "pdf_path": {"type": "string", "description": "Caminho absoluto do PDF de origem."},
                "text": {"type": "string", "description": "Texto da marca d'água (ex: 'CONFIDENCIAL')."},
                "output_name": {
                    "type": "string",
                    "description": "Nome do arquivo de saída sem extensão (padrão: 'watermarked').",
                    "default": "watermarked",
                },
            },
            "required": ["pdf_path", "text"],
        },
    },
]


def _safe_path(path: str) -> str:
    """Resolve and validate that path is under /opt/data."""
    real = os.path.realpath(path)
    if not real.startswith("/opt/data/"):
        raise ValueError(f"Path must be under /opt/data (got: {path})")
    return real


def _open_pdf(real: str):
    """pymupdf.open() with a friendly error instead of a raw parser exception —
    covers genuinely corrupted/malformed files (bad header, truncated upload,
    wrong-format-with-.pdf-extension). Does NOT special-case password-protected
    PDFs (doc.needs_pass) — open() itself succeeds for those; that's a distinct,
    not-yet-handled case (page access fails later with its own generic error)."""
    try:
        import pymupdf
    except ImportError:
        raise RuntimeError("pymupdf is not installed in this environment")
    try:
        return pymupdf.open(real)
    except Exception as e:
        raise ValueError(f"Não foi possível abrir o PDF — pode estar corrompido ou em formato inválido ({e}).")


def extract_pdf_text(
    pdf_path: str,
    start_page: int = None,
    end_page: int = None,
    max_chars: int = DEFAULT_MAX_TEXT_CHARS,
    max_ocr_pages: int = DEFAULT_MAX_OCR_PAGES,
):
    real = _safe_path(pdf_path)
    if not os.path.isfile(real):
        raise FileNotFoundError(f"PDF not found: {pdf_path}")

    doc = _open_pdf(real)
    total_pages = doc.page_count
    lo = max(1, int(start_page) if start_page else 1)
    hi = min(total_pages, int(end_page) if end_page else total_pages)
    if lo > hi:
        doc.close()
        raise ValueError(f"Intervalo de páginas inválido: {lo}-{hi} (o PDF tem {total_pages} página(s)).")

    page_texts = []
    ocr_pages = 0
    ocr_pages_skipped = 0

    for page_num in range(lo - 1, hi):
        page = doc[page_num]
        try:
            text = page.get_text().strip()
        except Exception as e:
            log(f"  get_text failed on page {page_num + 1}: {e}")
            text = ""
        if not text:
            # Page has no text layer (scanned/image-only) — fall back to OCR,
            # same logic the ocr-and-documents skill documents for execute_code,
            # capped so one bulky scanned document can't stall the whole call.
            if ocr_pages < max_ocr_pages:
                try:
                    tp = page.get_textpage_ocr(flags=0, language="por+eng", dpi=300, full=True)
                    text = page.get_text(textpage=tp).strip()
                    if text:
                        ocr_pages += 1
                except Exception as e:
                    log(f"  OCR failed on page {page_num + 1}: {e}")
            else:
                ocr_pages_skipped += 1
        if text:
            page_texts.append(f"--- Página {page_num + 1} ---\n{text}")

    doc.close()
    log(
        f"Text extracted from {real} (pages {lo}-{hi}/{total_pages}): "
        f"{len(page_texts)} with text, {ocr_pages} via OCR, {ocr_pages_skipped} OCR-skipped"
    )

    full_text = "\n\n".join(page_texts)
    truncated = len(full_text) > max_chars
    if truncated:
        full_text = full_text[:max_chars]

    return {
        "text": full_text,
        "total_pages": total_pages,
        "range": (lo, hi),
        "pages_with_text": len(page_texts),
        "ocr_pages": ocr_pages,
        "ocr_pages_skipped": ocr_pages_skipped,
        "truncated": truncated,
    }


def extract_pdf_images(pdf_path: str, max_images: int = 10, min_size: int = 100):
    real = _safe_path(pdf_path)
    if not os.path.isfile(real):
        raise FileNotFoundError(f"PDF not found: {pdf_path}")

    img_dir = os.path.join(OUTPUT_DIR, "pdf_images")
    os.makedirs(img_dir, exist_ok=True)

    doc = _open_pdf(real)
    extracted = []
    seen_xrefs = set()

    for page_num, page in enumerate(doc):
        if len(extracted) >= max_images:
            break
        for img_index, img in enumerate(page.get_images(full=True)):
            if len(extracted) >= max_images:
                break
            xref = img[0]
            if xref in seen_xrefs:
                continue
            seen_xrefs.add(xref)

            base_image = doc.extract_image(xref)
            width = base_image["width"]
            height = base_image["height"]
            if width < min_size or height < min_size:
                continue

            ext = base_image["ext"]
            filename = f"page{page_num + 1}_img{img_index + 1}.{ext}"
            filepath = os.path.join(img_dir, filename)
            with open(filepath, "wb") as f:
                f.write(base_image["image"])

            extracted.append({
                "path": filepath,
                "page": page_num + 1,
                "width": width,
                "height": height,
                "ext": ext,
            })
            log(f"  Extracted: {filename} ({width}x{height})")

    doc.close()
    log(f"Total extracted: {len(extracted)} images from {real}")
    return extracted


def extract_docx_text(docx_path: str, offset: int = 0, max_chars: int = DEFAULT_MAX_TEXT_CHARS):
    real = _safe_path(docx_path)
    if not os.path.isfile(real):
        raise FileNotFoundError(f"Documento não encontrado: {docx_path}")

    try:
        import docx
    except ImportError:
        raise RuntimeError("python-docx is not installed in this environment")

    try:
        document = docx.Document(real)
    except Exception as e:
        raise ValueError(f"Não foi possível abrir o documento Word — pode estar corrompido ou em formato inválido ({e}).")

    # python-docx reads paragraphs and tables as separate collections (no single
    # "body order" API without walking low-level XML), so a table embedded between
    # two paragraphs is appended after all paragraph text rather than inline. Fine
    # for "read and answer questions about this document" — order rarely matters
    # there, and it keeps this simple instead of hand-parsing document.element.body.
    parts = []
    for para in document.paragraphs:
        if para.text.strip():
            parts.append(para.text)
    for table_idx, table in enumerate(document.tables):
        rows = []
        for row in table.rows:
            cells = [cell.text.strip() for cell in row.cells]
            rows.append(" | ".join(cells))
        if rows:
            parts.append(f"[Tabela {table_idx + 1}]\n" + "\n".join(rows))

    full_text = "\n\n".join(parts)
    total_chars = len(full_text)
    offset = max(0, int(offset))
    window = full_text[offset:offset + max_chars]
    truncated = offset + len(window) < total_chars
    next_offset = offset + len(window) if truncated else None

    log(f"Text extracted from {real}: {total_chars} chars total, window [{offset}:{offset + len(window)}]")

    return {
        "text": window,
        "total_chars": total_chars,
        "offset": offset,
        "truncated": truncated,
        "next_offset": next_offset,
    }


def extract_xlsx_text(
    xlsx_path: str,
    sheet_name: str = None,
    max_chars: int = DEFAULT_MAX_TEXT_CHARS,
    max_rows_per_sheet: int = DEFAULT_MAX_ROWS_PER_SHEET,
):
    real = _safe_path(xlsx_path)
    if not os.path.isfile(real):
        raise FileNotFoundError(f"Planilha não encontrada: {xlsx_path}")

    try:
        import openpyxl
    except ImportError:
        raise RuntimeError("openpyxl is not installed in this environment")

    try:
        wb = openpyxl.load_workbook(real, data_only=True, read_only=True)
    except Exception as e:
        raise ValueError(f"Não foi possível abrir a planilha — pode estar corrompida ou em formato inválido ({e}).")

    sheet_names = wb.sheetnames
    target_sheets = [sheet_name] if sheet_name else sheet_names
    invalid = [s for s in target_sheets if s not in sheet_names]
    if invalid:
        raise ValueError(f"Aba(s) não encontrada(s): {', '.join(invalid)}. Abas disponíveis: {', '.join(sheet_names)}")

    parts = []
    rows_truncated_sheets = []
    for name in target_sheets:
        ws = wb[name]
        rows_text = []
        for i, row in enumerate(ws.iter_rows(values_only=True)):
            if i >= max_rows_per_sheet:
                rows_truncated_sheets.append(name)
                break
            cells = ["" if v is None else str(v) for v in row]
            if any(c.strip() for c in cells):
                rows_text.append(",".join(cells))
        parts.append(f"[Aba: {name}]\n" + "\n".join(rows_text))

    full_text = "\n\n".join(parts)
    chars_truncated = len(full_text) > max_chars
    if chars_truncated:
        full_text = full_text[:max_chars]

    log(f"Text extracted from {real}: sheets={sheet_names}, rows_truncated_in={rows_truncated_sheets}")

    return {
        "text": full_text,
        "sheet_names": sheet_names,
        "rows_truncated_sheets": rows_truncated_sheets,
        "chars_truncated": chars_truncated,
    }


def extract_pptx_text(pptx_path: str, start_slide: int = None, end_slide: int = None, max_chars: int = DEFAULT_MAX_TEXT_CHARS):
    real = _safe_path(pptx_path)
    if not os.path.isfile(real):
        raise FileNotFoundError(f"Apresentação não encontrada: {pptx_path}")

    try:
        import pptx as pptx_lib
    except ImportError:
        raise RuntimeError("python-pptx is not installed in this environment")

    try:
        prs = pptx_lib.Presentation(real)
    except Exception as e:
        raise ValueError(f"Não foi possível abrir a apresentação — pode estar corrompida ou em formato inválido ({e}).")

    total_slides = len(prs.slides)
    lo = max(1, int(start_slide) if start_slide else 1)
    hi = min(total_slides, int(end_slide) if end_slide else total_slides)
    if lo > hi:
        raise ValueError(f"Intervalo de slides inválido: {lo}-{hi} (a apresentação tem {total_slides} slide(s)).")

    parts = []
    for idx, slide in enumerate(prs.slides, start=1):
        if idx < lo or idx > hi:
            continue
        texts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                t = shape.text_frame.text.strip()
                if t:
                    texts.append(t)
            if shape.has_table:
                for row in shape.table.rows:
                    cells = [cell.text.strip() for cell in row.cells]
                    texts.append(" | ".join(cells))
        notes = ""
        if slide.has_notes_slide and slide.notes_slide.notes_text_frame.text.strip():
            notes = f"\n[Notas: {slide.notes_slide.notes_text_frame.text.strip()}]"
        if texts or notes:
            parts.append(f"--- Slide {idx} ---\n" + "\n".join(texts) + notes)

    full_text = "\n\n".join(parts)
    truncated = len(full_text) > max_chars
    if truncated:
        full_text = full_text[:max_chars]

    log(f"Text extracted from {real}: slides {lo}-{hi}/{total_slides}")

    return {
        "text": full_text,
        "total_slides": total_slides,
        "range": (lo, hi),
        "truncated": truncated,
    }


def read_text_file(file_path: str, max_chars: int = DEFAULT_MAX_TEXT_CHARS):
    real = _safe_path(file_path)
    if not os.path.isfile(real):
        raise FileNotFoundError(f"Arquivo não encontrado: {file_path}")

    ext = os.path.splitext(real)[1].lower()
    if ext not in TEXT_FILE_EXTENSIONS:
        raise ValueError(
            f"Extensão '{ext}' não é texto puro reconhecido — use a ferramenta específica "
            "(extract_pdf_text / extract_docx_text / extract_xlsx_text / extract_pptx_text) "
            "ou vision para imagens."
        )

    size = os.path.getsize(real)
    if size > MAX_TEXT_FILE_BYTES:
        raise ValueError(
            f"Arquivo muito grande ({size // 1024}KB) para leitura direta — "
            f"máximo {MAX_TEXT_FILE_BYTES // (1024 * 1024)}MB."
        )

    try:
        with open(real, "r", encoding="utf-8", errors="replace") as f:
            content = f.read()
    except Exception as e:
        raise ValueError(f"Não foi possível ler o arquivo como texto: {e}")

    truncated = len(content) > max_chars
    if truncated:
        content = content[:max_chars]

    log(f"Text read from {real}: {size} bytes")

    return {"text": content, "total_bytes": size, "truncated": truncated}


def _downscale_for_embed(path: str) -> str:
    """Return a path to a copy of the image capped at MAX_IMAGE_DIMENSION on its
    longer side, so base64-embedding many full-resolution photos (e.g. inspection
    report photos straight off a phone camera) doesn't bloat the HTML and stall
    chromium's render. Falls back to the original path if Pillow isn't available
    or the image can't be read — never blocks report generation on this step."""
    try:
        from PIL import Image
    except ImportError:
        return path

    try:
        with Image.open(path) as im:
            if max(im.size) <= MAX_IMAGE_DIMENSION:
                return path
            im.thumbnail((MAX_IMAGE_DIMENSION, MAX_IMAGE_DIMENSION))
            scaled_path = f"{path}.scaled.jpg"
            rgb = im.convert("RGB") if im.mode not in ("RGB", "L") else im
            rgb.save(scaled_path, "JPEG", quality=85)
            return scaled_path
    except Exception as e:
        log(f"  Downscale skipped for {path}: {e}")
        return path


def generate_pdf_report(title: str, images: list, subtitle: str = "", body: str = "",
                        output_name: str = "relatorio", images_per_row: int = 1):
    from datetime import date
    from html import escape

    today = date.today().strftime("%d/%m/%Y")

    # Validate all image paths
    valid_images = []
    for item in images:
        p = item.get("path", "")
        caption = item.get("caption", "")
        if not p:
            continue
        try:
            real = _safe_path(p)
        except ValueError as e:
            log(f"  Skipping image {p}: {e}")
            continue
        if not os.path.isfile(real):
            log(f"  Skipping missing image: {real}")
            continue
        valid_images.append({"path": real, "caption": caption})

    if not valid_images and not body.strip():
        raise ValueError(
            "No valid images and no body text found. Provide at least one valid image "
            "(path under /opt/data) or a 'body' with narrative text."
        )

    # Build HTML with embedded images as base64 to avoid chromium file:// security restrictions
    import base64
    import mimetypes

    def img_to_data_url(path: str) -> str:
        embed_path = _downscale_for_embed(path)
        mime, _ = mimetypes.guess_type(embed_path)
        if not mime:
            mime = "image/jpeg"
        with open(embed_path, "rb") as f:
            b64 = base64.b64encode(f.read()).decode("ascii")
        return f"data:{mime};base64,{b64}"

    col_width = "100%" if images_per_row == 1 else "48%"
    col_style = "display:inline-block;vertical-align:top;"

    img_html_parts = []
    for i, item in enumerate(valid_images):
        try:
            data_url = img_to_data_url(item["path"])
        except Exception as e:
            log(f"  Failed to encode image {item['path']}: {e}")
            continue
        caption_html = f'<p class="caption">{escape(item["caption"])}</p>' if item["caption"] else ""
        img_html_parts.append(
            f'<div style="width:{col_width};{col_style}margin:8px 1%;">'
            f'<img src="{data_url}" style="max-width:100%;height:auto;border:1px solid #ddd;">'
            f'{caption_html}</div>'
        )

    subtitle_html = f'<p class="subtitle">{escape(subtitle)}</p>' if subtitle else ""
    images_html = "\n".join(img_html_parts)

    # body is plain text (no HTML from the model) — escape, then turn blank-line-
    # separated blocks into paragraphs so multi-paragraph narrative renders cleanly.
    body_html = ""
    if body.strip():
        paragraphs = [p.strip() for p in body.strip().split("\n\n") if p.strip()]
        body_html = "\n".join(f"<p>{escape(p).replace(chr(10), '<br>')}</p>" for p in paragraphs)
        body_html = f'<div class="body">{body_html}</div>'

    html = f"""<!DOCTYPE html>
<html lang="pt-BR">
<head>
<meta charset="UTF-8">
<style>
  body {{ font-family: Arial, sans-serif; margin: 30px; font-size: 13px; color: #222; }}
  h1 {{ color: #1a1a2e; border-bottom: 2px solid #6c3483; padding-bottom: 8px; margin-bottom: 4px; }}
  .subtitle {{ color: #555; margin: 4px 0 8px 0; font-size: 12px; }}
  .date {{ color: #777; font-size: 11px; margin-bottom: 20px; }}
  .body {{ margin-bottom: 20px; line-height: 1.5; }}
  .body p {{ margin: 0 0 10px 0; }}
  .caption {{ text-align: center; color: #444; font-size: 11px; margin-top: 4px; }}
  .images-grid {{ width: 100%; }}
</style>
</head>
<body>
<h1>{escape(title)}</h1>
{subtitle_html}
<p class="date">Data: {today}</p>
{body_html}
<div class="images-grid">
{images_html}
</div>
</body>
</html>"""

    html_path = os.path.join(OUTPUT_DIR, f"{output_name}.html")
    pdf_path = os.path.join(OUTPUT_DIR, f"{output_name}.pdf")

    os.makedirs(OUTPUT_DIR, exist_ok=True)
    with open(html_path, "w", encoding="utf-8") as f:
        f.write(html)
    log(f"HTML written: {html_path} ({len(html)//1024}KB)")

    # Convert to PDF via chromium
    import subprocess
    cmd = [
        "/usr/bin/chromium",
        "--headless",
        "--no-sandbox",
        "--disable-dev-shm-usage",
        "--disable-gpu",
        "--run-all-compositor-stages-before-draw",
        "--print-to-pdf-no-header",
        f"--print-to-pdf={pdf_path}",
        f"file://{html_path}",
    ]
    try:
        result = subprocess.run(cmd, capture_output=True, timeout=CHROMIUM_TIMEOUT_SECONDS)
    except subprocess.TimeoutExpired:
        raise RuntimeError(
            f"Chromium timed out after {CHROMIUM_TIMEOUT_SECONDS}s — "
            f"try fewer images or a lower images_per_row"
        )
    if result.returncode != 0:
        err = result.stderr.decode("utf-8", errors="replace")[:300]
        raise RuntimeError(f"Chromium failed (exit {result.returncode}): {err}")

    if not os.path.isfile(pdf_path):
        raise RuntimeError(f"Chromium ran but PDF was not created at {pdf_path}")

    size_kb = os.path.getsize(pdf_path) // 1024
    log(f"PDF generated: {pdf_path} ({size_kb}KB)")
    return pdf_path


def pdf_split(pdf_path: str, start_page: int, end_page: int, output_name: str = "split"):
    real = _safe_path(pdf_path)
    if not os.path.isfile(real):
        raise FileNotFoundError(f"PDF not found: {pdf_path}")

    import pymupdf

    doc = _open_pdf(real)
    total_pages = doc.page_count
    lo = max(1, int(start_page))
    hi = min(total_pages, int(end_page))
    if lo > hi:
        doc.close()
        raise ValueError(f"Intervalo de páginas inválido: {lo}-{hi} (o PDF tem {total_pages} página(s)).")

    new_doc = pymupdf.open()
    new_doc.insert_pdf(doc, from_page=lo - 1, to_page=hi - 1)
    doc.close()

    os.makedirs(OUTPUT_DIR, exist_ok=True)
    out_path = os.path.join(OUTPUT_DIR, f"{output_name}.pdf")
    new_doc.save(out_path)
    new_doc.close()
    log(f"Split {real} pages {lo}-{hi} -> {out_path}")
    return out_path, lo, hi, total_pages


def pdf_merge(pdf_paths: list, output_name: str = "merged"):
    if not pdf_paths or len(pdf_paths) < 2:
        raise ValueError("Forneça pelo menos 2 caminhos de PDF para mesclar.")

    import pymupdf

    result = pymupdf.open()
    page_counts = []
    for p in pdf_paths:
        real = _safe_path(p)
        if not os.path.isfile(real):
            result.close()
            raise FileNotFoundError(f"PDF não encontrado: {p}")
        src = _open_pdf(real)
        result.insert_pdf(src)
        page_counts.append(src.page_count)
        src.close()

    os.makedirs(OUTPUT_DIR, exist_ok=True)
    out_path = os.path.join(OUTPUT_DIR, f"{output_name}.pdf")
    result.save(out_path)
    total_pages = result.page_count
    result.close()
    log(f"Merged {len(pdf_paths)} PDFs ({page_counts} pages each) -> {out_path} ({total_pages} pages)")
    return out_path, total_pages


def pdf_watermark(pdf_path: str, text: str, output_name: str = "watermarked"):
    real = _safe_path(pdf_path)
    if not os.path.isfile(real):
        raise FileNotFoundError(f"PDF not found: {pdf_path}")
    if not text.strip():
        raise ValueError("text is required")

    import pymupdf

    doc = _open_pdf(real)
    fontsize = 48
    text_length = pymupdf.get_text_length(text, fontsize=fontsize)

    for page in doc:
        rect = page.rect
        center = pymupdf.Point(rect.width / 2, rect.height / 2)
        # insert_text's own `rotate` only accepts multiples of 90 — `morph` (a
        # pivot point + rotation matrix) is what actually gives a diagonal stamp.
        morph = (center, pymupdf.Matrix(45))
        page.insert_text(
            (rect.width / 2 - text_length / 2, rect.height / 2),
            text,
            fontsize=fontsize,
            color=(0.8, 0, 0),
            fill_opacity=0.3,
            morph=morph,
        )

    os.makedirs(OUTPUT_DIR, exist_ok=True)
    out_path = os.path.join(OUTPUT_DIR, f"{output_name}.pdf")
    doc.save(out_path)
    total_pages = doc.page_count
    doc.close()
    log(f"Watermarked {real} ({total_pages} pages) -> {out_path}")
    return out_path, total_pages


def call_tool(name: str, args: dict):
    if name == "extract_pdf_text":
        pdf_path = args.get("pdf_path", "")
        if not pdf_path:
            raise ValueError("pdf_path is required")

        start_page = args.get("start_page")
        end_page = args.get("end_page")
        max_chars = int(args.get("max_chars") or DEFAULT_MAX_TEXT_CHARS)
        max_ocr_pages = int(args.get("max_ocr_pages") or DEFAULT_MAX_OCR_PAGES)
        result = extract_pdf_text(
            pdf_path, start_page=start_page, end_page=end_page,
            max_chars=max_chars, max_ocr_pages=max_ocr_pages,
        )
        lo, hi = result["range"]

        if not result["text"].strip():
            return (
                f"Não encontrei texto legível nas páginas {lo}-{hi} deste PDF "
                f"({result['total_pages']} página(s) no total) — nem no texto nativo nem via OCR. "
                "Pode estar protegido, corrompido, ou ser um escaneado de qualidade muito baixa."
            )

        header = f"Texto extraído das páginas {lo}-{hi}/{result['total_pages']}"
        if result["ocr_pages"]:
            header += f" ({result['ocr_pages']} via OCR)"
        header += ":\n\n"

        footer_bits = []
        if result["truncated"]:
            footer_bits.append(
                f"texto truncado em {max_chars} caracteres — peça um max_chars maior, ou um "
                "start_page/end_page mais estreito, se precisar do restante"
            )
        if result["ocr_pages_skipped"]:
            footer_bits.append(
                f"{result['ocr_pages_skipped']} página(s) escaneada(s) além do limite de OCR "
                f"({max_ocr_pages}/chamada) não foram lidas — peça um start_page/end_page mais "
                "estreito para alcançá-las"
            )
        footer = f"\n\n[{'; '.join(footer_bits)}]" if footer_bits else ""
        return header + result["text"] + footer

    elif name == "extract_pdf_images":
        pdf_path = args.get("pdf_path", "")
        if not pdf_path:
            raise ValueError("pdf_path is required")

        max_images = int(args.get("max_images") or 10)
        min_size = int(args.get("min_size") or 100)
        images = extract_pdf_images(pdf_path, max_images=max_images, min_size=min_size)

        if not images:
            return (
                f"Nenhuma imagem embutida encontrada neste PDF "
                f"(filtro de tamanho mínimo: {min_size}px). "
                "O PDF pode ser escaneado (imagem da página inteira) ou não ter fotos embutidas."
            )

        lines = [f"Extraídas {len(images)} imagens:"]
        for img in images:
            lines.append(f"  página {img['page']}: {img['width']}x{img['height']}px → {img['path']}")
        lines.append("")
        lines.append("Para enviar cada imagem ao chat, inclua uma linha MEDIA:<path> na sua resposta:")
        for img in images:
            lines.append(f"MEDIA:{img['path']}")
        return "\n".join(lines)

    elif name == "extract_docx_text":
        docx_path = args.get("docx_path", "")
        if not docx_path:
            raise ValueError("docx_path is required")

        offset = int(args.get("offset") or 0)
        max_chars = int(args.get("max_chars") or DEFAULT_MAX_TEXT_CHARS)
        result = extract_docx_text(docx_path, offset=offset, max_chars=max_chars)

        if not result["text"].strip():
            return "Não encontrei texto neste documento Word (pode estar vazio ou só ter imagens)."

        header = f"Texto extraído (caracteres {result['offset']}-{result['offset'] + len(result['text'])} de {result['total_chars']}):\n\n"
        footer = ""
        if result["truncated"]:
            footer = f"\n\n[truncado — para continuar, chame de novo com offset={result['next_offset']}]"
        return header + result["text"] + footer

    elif name == "extract_xlsx_text":
        xlsx_path = args.get("xlsx_path", "")
        if not xlsx_path:
            raise ValueError("xlsx_path is required")

        sheet_name = args.get("sheet_name") or None
        max_chars = int(args.get("max_chars") or DEFAULT_MAX_TEXT_CHARS)
        max_rows_per_sheet = int(args.get("max_rows_per_sheet") or DEFAULT_MAX_ROWS_PER_SHEET)
        result = extract_xlsx_text(
            xlsx_path, sheet_name=sheet_name, max_chars=max_chars, max_rows_per_sheet=max_rows_per_sheet,
        )

        if not result["text"].strip():
            return f"Planilha vazia. Abas encontradas: {', '.join(result['sheet_names'])}."

        header = f"Conteúdo de {', '.join(result['sheet_names']) if not sheet_name else sheet_name}:\n\n"
        footer_bits = []
        if result["rows_truncated_sheets"]:
            footer_bits.append(
                f"linhas cortadas em {max_rows_per_sheet}/aba nas abas: "
                f"{', '.join(result['rows_truncated_sheets'])} — chame de novo com sheet_name para ler mais"
            )
        if result["chars_truncated"]:
            footer_bits.append(f"texto truncado em {max_chars} caracteres")
        footer = f"\n\n[{'; '.join(footer_bits)}]" if footer_bits else ""
        return header + result["text"] + footer

    elif name == "extract_pptx_text":
        pptx_path = args.get("pptx_path", "")
        if not pptx_path:
            raise ValueError("pptx_path is required")

        start_slide = args.get("start_slide")
        end_slide = args.get("end_slide")
        max_chars = int(args.get("max_chars") or DEFAULT_MAX_TEXT_CHARS)
        result = extract_pptx_text(pptx_path, start_slide=start_slide, end_slide=end_slide, max_chars=max_chars)
        lo, hi = result["range"]

        if not result["text"].strip():
            return f"Não encontrei texto nos slides {lo}-{hi} ({result['total_slides']} slide(s) no total)."

        header = f"Texto extraído dos slides {lo}-{hi}/{result['total_slides']}:\n\n"
        footer = ""
        if result["truncated"]:
            footer = f"\n\n[texto truncado em {max_chars} caracteres — peça um start_slide/end_slide mais estreito se precisar do restante]"
        return header + result["text"] + footer

    elif name == "read_text_file":
        file_path = args.get("file_path", "")
        if not file_path:
            raise ValueError("file_path is required")

        max_chars = int(args.get("max_chars") or DEFAULT_MAX_TEXT_CHARS)
        result = read_text_file(file_path, max_chars=max_chars)

        if not result["text"].strip():
            return "Arquivo vazio."

        footer = f"\n\n[truncado em {max_chars} caracteres — peça um max_chars maior se precisar do restante]" if result["truncated"] else ""
        return result["text"] + footer

    elif name == "generate_pdf_report":
        title = args.get("title", "")
        if not title:
            raise ValueError("title is required")
        images = args.get("images", [])
        body = args.get("body", "") or ""

        subtitle = args.get("subtitle", "")
        output_name = args.get("output_name", "relatorio") or "relatorio"
        images_per_row = int(args.get("images_per_row") or 1)

        pdf_path = generate_pdf_report(
            title=title,
            images=images,
            subtitle=subtitle,
            body=body,
            output_name=output_name,
            images_per_row=images_per_row,
        )
        size_kb = os.path.getsize(pdf_path) // 1024
        return (
            f"PDF gerado com sucesso: {pdf_path} ({size_kb}KB)\n\n"
            f"Para enviar o PDF ao usuário, inclua na sua resposta:\n"
            f"MEDIA:{pdf_path}"
        )

    elif name == "pdf_split":
        pdf_path = args.get("pdf_path", "")
        if not pdf_path:
            raise ValueError("pdf_path is required")
        start_page = args.get("start_page")
        end_page = args.get("end_page")
        if start_page is None or end_page is None:
            raise ValueError("start_page and end_page are required")
        output_name = args.get("output_name", "split") or "split"

        out_path, lo, hi, total_pages = pdf_split(pdf_path, start_page, end_page, output_name=output_name)
        size_kb = os.path.getsize(out_path) // 1024
        return (
            f"PDF com as páginas {lo}-{hi} (de {total_pages}) gerado: {out_path} ({size_kb}KB)\n\n"
            f"Para enviar ao usuário, inclua na sua resposta:\nMEDIA:{out_path}"
        )

    elif name == "pdf_merge":
        pdf_paths = args.get("pdf_paths", [])
        output_name = args.get("output_name", "merged") or "merged"

        out_path, total_pages = pdf_merge(pdf_paths, output_name=output_name)
        size_kb = os.path.getsize(out_path) // 1024
        return (
            f"PDF mesclado ({len(pdf_paths)} arquivos, {total_pages} páginas) gerado: {out_path} ({size_kb}KB)\n\n"
            f"Para enviar ao usuário, inclua na sua resposta:\nMEDIA:{out_path}"
        )

    elif name == "pdf_watermark":
        pdf_path = args.get("pdf_path", "")
        if not pdf_path:
            raise ValueError("pdf_path is required")
        text = args.get("text", "")
        if not text:
            raise ValueError("text is required")
        output_name = args.get("output_name", "watermarked") or "watermarked"

        out_path, total_pages = pdf_watermark(pdf_path, text, output_name=output_name)
        size_kb = os.path.getsize(out_path) // 1024
        return (
            f"PDF com marca d'água '{text}' aplicada em {total_pages} página(s): {out_path} ({size_kb}KB)\n\n"
            f"Para enviar ao usuário, inclua na sua resposta:\nMEDIA:{out_path}"
        )

    else:
        raise ValueError(f"Unknown tool: {name}")


def handle(msg: dict):
    id_ = msg.get("id")
    method = msg.get("method", "")
    params = msg.get("params") or {}

    if method == "initialize":
        return reply(id_, {
            "protocolVersion": PROTOCOL_VERSION,
            "capabilities": {"tools": {}},
            "serverInfo": {"name": SERVER_NAME, "version": SERVER_VERSION},
        })

    if method == "notifications/initialized":
        return

    if method == "tools/list":
        return reply(id_, {"tools": TOOLS})

    if method == "tools/call":
        tool_name = params.get("name", "")
        tool_args = params.get("arguments") or {}
        try:
            text = call_tool(tool_name, tool_args)
            return reply(id_, {"content": [{"type": "text", "text": text}]})
        except Exception as e:
            msg_text = str(e)
            log(f"Error in {tool_name}: {msg_text}")
            return reply(id_, {
                "content": [{"type": "text", "text": f"Erro: {msg_text}"}],
                "isError": True,
            })

    if id_ is not None:
        reply_error(id_, -32601, f"Method not found: {method}")


def main():
    log(f"Started ({SERVER_NAME} {SERVER_VERSION})")
    for line in sys.stdin:
        line = line.strip()
        if not line:
            continue
        try:
            msg = json.loads(line)
        except json.JSONDecodeError:
            log(f"Invalid JSON: {line[:80]}")
            continue
        try:
            handle(msg)
        except Exception as e:
            log(f"Unhandled error: {e}")
            if msg.get("id") is not None:
                reply_error(msg["id"], -32603, "Internal error")


if __name__ == "__main__":
    main()
